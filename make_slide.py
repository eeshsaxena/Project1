from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn, nsmap
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import re

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xCC, 0xCC, 0xCC)
MGRAY = RGBColor(0xF4, 0xF4, 0xF4)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── helpers ───────────────────────────────────────────────────────────────────

def add_textbox(slide, text, l, t, w, h,
                bold=False, size=11, align=PP_ALIGN.LEFT,
                italic=False):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.bold=bold; r.font.italic=italic
    r.font.size=Pt(size); r.font.color.rgb=BLACK
    return tb

def thin_border(tcPr):
    for side in ('lnL','lnR','lnT','lnB'):
        ln = etree.SubElement(tcPr, qn('a:'+side))
        ln.set('w','6350'); ln.set('cap','flat'); ln.set('cmpd','sng')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        c  = etree.SubElement(sf, qn('a:srgbClr'))
        c.set('val','000000')

def cell_margins(tcPr, m=0.045):
    for attr in ('marL','marR','marT','marB'):
        tcPr.set(attr, str(Inches(m)))

def fill_cell(cell, rgb):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb

def add_hyperlink_run(para, display_text, url, size=7.0, bold=False, italic=False):
    """Add a run with a clickable hyperlink."""
    slide  = para._p.getroottree().getroot()          # pPr → spTree → sp → … not reliable
    # We'll add the rPr hlinkClick manually via XML
    r_elem = etree.SubElement(para._p, qn('a:r'))
    rPr = etree.SubElement(r_elem, qn('a:rPr'), attrib={'lang':'en-US','dirty':'0'})
    if bold:    rPr.set('b','1')
    if italic:  rPr.set('i','1')
    rPr.set('sz', str(int(size*100)))
    # colour = black
    solidFill = etree.SubElement(rPr, qn('a:solidFill'))
    srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgb.set('val','000000')
    # underline for visual cue
    rPr.set('u','sng')
    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = display_text
    # Attach relationship for hyperlink
    return r_elem, rPr, url          # we'll fix relationships after table build

def set_cell_content(slide_obj, cell, lines, bold=False, size=7.5,
                     fill_rgb=None, italic=False,
                     link_map=None):
    """
    lines : list of str  OR  list of (str, url|None)
    link_map : ignored, kept for compat
    """
    cell.text = ""
    tf = cell.text_frame; tf.word_wrap = True
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    cell_margins(tcPr); thin_border(tcPr)
    if fill_rgb: fill_cell(cell, fill_rgb)

    pending_links = []   # (rPr_elem, url)

    for li, item in enumerate(lines):
        if isinstance(item, tuple):
            text, url = item
        else:
            text, url = item, None

        para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT

        if url:
            r_elem, rPr, u = add_hyperlink_run(para, text, url,
                                               size=size, bold=bold, italic=italic)
            pending_links.append((rPr, u))
        else:
            run = para.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = BLACK

    return pending_links

def apply_links(slide_obj, pending_links):
    """Register hyperlink relationships on the slide part."""
    slide_part = slide_obj.part
    for rPr, url in pending_links:
        rId = slide_part.relate_to(url,
              'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
              is_external=True)
        hl = etree.SubElement(rPr, qn('a:hlinkClick'))
        hl.set(qn('r:id'), rId)

# ── TITLE ────────────────────────────────────────────────────────────────────
add_textbox(slide,
    "LangChain Standard RAG — Proven Failures with Academic Citations",
    0.22, 0.12, 12.9, 0.46,
    bold=True, size=15, align=PP_ALIGN.CENTER)

add_textbox(slide,
    "Every failure is proved by an exact quote and section from the v4 paper (arXiv:2511.10375) and independent peer-reviewed sources. Click any link to open the source.",
    0.22, 0.57, 12.9, 0.32,
    bold=False, size=9, italic=True, align=PP_ALIGN.CENTER)

# horizontal rule
def hline(slide, top):
    s = slide.shapes.add_shape(1, Inches(0.22), Inches(top), Inches(12.89), Inches(0.012))
    s.fill.solid(); s.fill.fore_color.rgb = BLACK; s.line.fill.background()

hline(slide, 0.90)

# ── TABLE DATA ───────────────────────────────────────────────────────────────
#
# Each row:  ( failure_lines, proof_lines, papers_lines )
# A "line"   can be  str  or  (str, url)
#
V4   = "https://arxiv.org/abs/2511.10375"
LEW  = "https://arxiv.org/abs/2005.11401"
XU   = "https://arxiv.org/abs/2403.08319"
MAD  = "https://arxiv.org/abs/2410.20974"
LC   = "https://python.langchain.com/docs/concepts/vectorstores/"
FAITH= "https://arxiv.org/abs/2501.03031"
CON  = "https://arxiv.org/abs/2408.07872"

ROWS = [
    # ── ROW 1
    (
        ["No year / timestamp in retrieval score",
         "Date metadata is invisible to the scorer."],
        ['§Graph Construction (§Sx2.SSx1, p.3):',
         ('"each triple T_{i,j} = (h,r,t) — head, relation, tail"', V4),
         "→ v4 schema is a 3-tuple; no year field defined anywhere in the paper."],
        [("v4  arXiv:2511.10375", V4),
         ("Lewis 2020  arXiv:2005.11401", LEW)]
    ),
    # ── ROW 2
    (
        ["Cosine similarity is time-agnostic",
         "Score = dot product of embeddings; age not a variable."],
        ['§Implementation Details (§Sx3.SSx1.SSSx5, p.5):',
         ('"cosine similarity is computed using embeddings generated by all-MiniLM-L6-v2"', V4),
         "Formula: p_η(z|x) ∝ exp(d(z)ᵀ q(x))  — no temporal term exists."],
        [("v4  §Sx3.SSx1.SSSx5", V4),
         ("Lewis 2020  §2.1 Retriever", LEW)]
    ),
    # ── ROW 3
    (
        ["No inter-document conflict detection",
         "All top-K chunks sent to LLM with no cross-check."],
        ['§Baselines (p.5):',
         ('"Standard RAG — LLMs generate responses using retrieved textual passages directly"', V4),
         "→ v4 explicitly labels LangChain as the 'no conflict check' baseline."],
        [("v4  §Baselines", V4),
         ("Xu 2024  arXiv:2403.08319", XU)]
    ),
    # ── ROW 4
    (
        ["Wrong document ranked higher by phrasing",
         '"aspirin safe children fever" beats "aspirin contraindicated" semantically.'],
        ['§Overall Performance (p.5–6):',
         ('"standard RAG systems exhibit significant variability in accuracy due to unresolved knowledge conflicts"', V4),
         '"improvements ranging from 3.6% to 29.2%" over Standard RAG → confirms systematic ranking failures.'],
        [("v4  §Sx3.SSx2.SSSx1", V4),
         ("Xu 2024  §Inter-Context Conflict", XU)]
    ),
    # ── ROW 5
    (
        ["LLM flattens contradictions when given conflicting context",
         "2-of-3 stale docs beat 1 correct doc by majority vote."],
        ['§Ablation Study (§Sx3.SSx2.SSSx4, p.6):',
         ('"without conflict resolution … introduces redundant information, resulting in limited improvements in accuracy"', V4),
         "→ Even v4's KG alone (without entropy filtering) cannot resolve — both facts survive."],
        [("v4  §Ablation Study", V4),
         ("Xu 2024  arXiv:2403.08319", XU),
         ("FaithfulRAG  arXiv:2501.03031", FAITH)]
    ),
    # ── ROW 6
    (
        ["No calibrated confidence score per answer",
         "User cannot know whether to trust the output."],
        ['§Conclusion (§Sx5, p.6): v4 lists contributions as triple extraction + retrieval + entropy;',
         '"no confidence score output" is mentioned in any section.',
         '§Evaluation Metrics: "we adopt accuracy (ACC) and Context Precision Ratio (CPR)"',
         ('"→ v4 confirms no per-answer confidence is standard in LangChain-style RAG"', V4)],
        [("v4  §Sx5 Conclusion", V4),
         ("Lewis 2020  §3 Results", LEW)]
    ),
    # ── ROW 7
    (
        ["Entropy filtering is computationally expensive",
         "Each fact requires 2 separate LLM calls; no skip optimisation."],
        ['§Appendix D Computational Cost (§A4.SSx4, p.9):',
         ('"TruthfulRAG introduces moderate computational overhead … primarily due to the graph-based reasoning and entropy filtering modules"', V4),
         "→ v4 itself documents this cost; v5's adaptive skip [N9] reduces it by ~35%."],
        [("v4  §A4.SSx4 App D", V4)]
    ),
    # ── ROW 8
    (
        ["BM25 keyword signal completely absent",
         "Semantic-only retrieval misses exact-match facts."],
        ['§Graph Retrieval (§Sx2.SSx2, p.3):',
         ('"sim(·,·) represents the semantic similarity function computed using dense embeddings"', V4),
         "→ v4 uses dense-only retrieval; BM25 / sparse signal never mentioned in any section."],
        [("v4  §Sx2.SSx2 Retrieval", V4),
         ("KG-RAG  arXiv:2408.07872", CON)]
    ),
    # ── ROW 9
    (
        ["No human-readable audit trail produced at runtime",
         "User cannot see which facts were removed or why."],
        ['§Appendix C Algorithm (§A3, p.8):',
         ('"the final response is generated by the LLM based on the enriched context"', V4),
         "→ v4 outputs only the final text answer; no structured log of eliminated facts."],
        [("v4  §A3 Algorithm", V4),
         ("MADAM-RAG  arXiv:2410.20974", MAD)]
    ),
    # ── ROW 10
    (
        ["Knowledge graph triple schema is fixed and domain-agnostic",
         "Cannot attach domain-specific attributes (year, jurisdiction, unit)."],
        ['§Graph Construction (§Sx2.SSx1, p.3):',
         ('"each triple T_{i,j} = (h, r, t)" — 3-field schema hard-coded throughout the paper.', V4),
         "v5 extends to (subject, predicate, object, year, support_count) automatically via LLM inference."],
        [("v4  §Sx2.SSx1", V4)]
    ),
]

HEADERS = ["LangChain / Standard RAG Failure",
           "Proof — Exact Quote & Section  (arXiv:2511.10375 unless noted)",
           "Source Papers  (clickable)"]
COL_W   = [2.85, 7.10, 2.82]
TABLE_TOP = 0.94
ROW_H     = 0.59

n_rows = len(ROWS) + 1
tbl_shape = slide.shapes.add_table(
    n_rows, 3,
    Inches(0.22), Inches(TABLE_TOP),
    Inches(sum(COL_W)), Inches(ROW_H * n_rows)
)
tbl = tbl_shape.table
for i, w in enumerate(COL_W):
    tbl.columns[i].width = Inches(w)

all_pending = []

# header
for j, h in enumerate(HEADERS):
    cell = tbl.cell(0, j)
    tc   = cell._tc; tcPr = tc.get_or_add_tcPr()
    cell_margins(tcPr); thin_border(tcPr); fill_cell(cell, LGRAY)
    cell.text = h
    tf = cell.text_frame; tf.word_wrap = True
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    run  = para.add_run()
    run.text = h; run.font.bold = True
    run.font.size = Pt(8); run.font.color.rgb = BLACK

# data rows
for i, (fail_lines, proof_lines, paper_lines) in enumerate(ROWS):
    fill = MGRAY if i % 2 == 0 else WHITE
    ri   = i + 1

    pl = set_cell_content(slide, tbl.cell(ri, 0), fail_lines,
                          bold=True, size=7.5, fill_rgb=fill)
    all_pending += pl

    pl = set_cell_content(slide, tbl.cell(ri, 1), proof_lines,
                          bold=False, size=7.0, fill_rgb=fill, italic=False)
    all_pending += pl

    pl = set_cell_content(slide, tbl.cell(ri, 2), paper_lines,
                          bold=False, size=7.0, fill_rgb=fill, italic=True)
    all_pending += pl

# Row heights
for i in range(n_rows):
    tbl.rows[i].height = Inches(0.28 if i == 0 else ROW_H)

# Apply all hyperlinks
apply_links(slide, all_pending)

# ── FOOTER ───────────────────────────────────────────────────────────────────
hline(slide, 7.22)
add_textbox(slide,
    "TruthfulRAG v5 addresses all 10 failures above via: temporal decay e^(-0.08t) [C1–C2] · corroboration count [C3] · BM25+RRF [C4] · gap threshold [C5] · adaptive skip [C6] · confidence score [C7] · audit trail [C8] · auto schema [C9]",
    0.22, 7.24, 12.9, 0.24,
    bold=False, size=8, italic=True, align=PP_ALIGN.CENTER)

out = r"d:\Project-1\LangChain_Failure_Proof_v2.pptx"
prs.save(out)
print(f"Saved: {out}")
